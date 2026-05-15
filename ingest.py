import os
import json
from pypdf import PdfReader
from pptx import Presentation

def process_files():
    raw_dir = "data/raw"
    processed_data = []

    # 檢查目錄是否存在
    if not os.path.exists(raw_dir):
        print(f"❌ 找不到目錄: {raw_dir}，請先執行 python make_data.py")
        return

    print("🚀 正在以穩定模式處理瑞鼎技術文件...")
    
    for filename in os.listdir(raw_dir):
        file_path = os.path.join(raw_dir, filename)
        content = ""
        
        try:
            if filename.endswith(".pdf"):
                print(f"正在讀取 PDF: {filename}")
                reader = PdfReader(file_path)
                content = " ".join([page.extract_text() for page in reader.pages if page.extract_text()])
            
            elif filename.endswith(".pptx"):
                print(f"正在讀取 PPTX: {filename}")
                prs = Presentation(file_path)
                text_runs = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_runs.append(shape.text)
                content = " ".join(text_runs)
                
            if content:
                processed_data.append({
                    "source": filename,
                    "content": content.strip(),
                    "category": "IC_Spec" if "spec" in filename.lower() else "Technical_Slide"
                })
        except Exception as e:
            print(f"❌ 處理 {filename} 時出錯: {e}")

    # 儲存結果
    os.makedirs("data/processed", exist_ok=True)
    output_path = "data/processed/knowledge_base.json"
    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(processed_data, f, ensure_ascii=False, indent=4)
    
    print(f"✅ 處理完成！結果已存至: {output_path}")

if __name__ == "__main__":
    process_files()