import os
from fpdf import FPDF
from pptx import Presentation

# 確保資料夾存在
os.makedirs("data/raw", exist_ok=True)

# 1. 產生模擬瑞鼎 IC 規格 PDF
pdf = FPDF()
pdf.add_page()
pdf.set_font("Arial", 'B', 16)
pdf.cell(40, 10, "Raydium Semiconductor - IC Spec Sheet")
pdf.ln(20)
pdf.set_font("Arial", size=12)
pdf.multi_cell(0, 10, "Project: RD-AI-2026 Inference Engine\n"
                     "Operating Voltage: 0.8V - 1.2V\n"
                     "Interface: MIPI DSI-2\n"
                     "Description: This chip provides high-efficiency AI upscaling for mobile displays.")
pdf.output("data/raw/specs.pdf")

# 2. 產生技術簡報 PPTX
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Raydium Display AI Integration"
slide.placeholders[1].text = "Optimizing driver performance with MCP-based LLM workflows."
prs.save("data/raw/presentation.pptx")

print("✅ 瑞鼎測試資料已產生於 data/raw/")